"""Собирает docs/презентация.pptx — стильная презентация для защиты курсовой.

Тёмная тема с акцентным синим, заголовки крупные, контент структурирован.
Картинки берутся из docs/diagrams/png/ и docs/screenshots/.

Запуск: python scripts/build_presentation_pptx.py
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu


ROOT = Path(__file__).resolve().parent.parent
DOCS = ROOT / "docs"
PNG = DOCS / "diagrams" / "png"
SCR = DOCS / "screenshots"
OUT = DOCS / "презентация.pptx"

# Цветовая схема (тёмная тема в стиле Material/Dracula)
BG_DARK = RGBColor(0x1A, 0x1F, 0x2E)
BG_PANEL = RGBColor(0x24, 0x2B, 0x3D)
TXT_PRIMARY = RGBColor(0xF2, 0xF4, 0xF8)
TXT_MUTED = RGBColor(0xA0, 0xA8, 0xBC)
ACCENT = RGBColor(0x29, 0xB6, 0xF6)  # светло-синий
ACCENT_2 = RGBColor(0x66, 0xBB, 0x6A)  # зелёный для метрик
ACCENT_3 = RGBColor(0xFF, 0xB0, 0x20)  # янтарный

# Размер слайда — 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_bg(slide, color: RGBColor = BG_DARK) -> None:
    """Полностью закрашивает фон слайда."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    # Отправить на задний план
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)


def add_text(
    slide, left, top, width, height,
    text: str,
    *,
    size: int = 18,
    color: RGBColor = TXT_PRIMARY,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "Segoe UI"
    return box


def add_bullets(slide, left, top, width, height, items: list[str], size: int = 16) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        # маркер цветной
        bullet = p.add_run()
        bullet.text = "▸  "
        bullet.font.size = Pt(size)
        bullet.font.color.rgb = ACCENT
        bullet.font.bold = True
        bullet.font.name = "Segoe UI"
        # текст
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = TXT_PRIMARY
        run.font.name = "Segoe UI"


def add_accent_line(slide, left, top, width: int = Inches(1)) -> None:
    """Тонкая цветная подчёркивающая полоска."""
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Emu(38000))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def slide_header(slide, number: int, title: str) -> None:
    """Стандартный header: номер + заголовок + подчёрк."""
    add_text(slide, Inches(0.5), Inches(0.35), Inches(0.8), Inches(0.5),
             f"{number:02d}", size=14, color=TXT_MUTED, bold=True)
    add_text(slide, Inches(0.5), Inches(0.7), Inches(12), Inches(0.9),
             title, size=32, color=TXT_PRIMARY, bold=True)
    add_accent_line(slide, Inches(0.5), Inches(1.55), Inches(1.2))


def make_blank_slide(prs: Presentation):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    return slide


# --------- Слайды ---------

def slide_title(prs: Presentation) -> None:
    slide = make_blank_slide(prs)

    # Большая декоративная полоса слева
    side = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.4), SLIDE_H)
    side.fill.solid()
    side.fill.fore_color.rgb = ACCENT
    side.line.fill.background()

    add_text(slide, Inches(1.2), Inches(1.5), Inches(11), Inches(0.6),
             "КУРСОВАЯ РАБОТА", size=18, color=ACCENT, bold=True)

    add_text(slide, Inches(1.2), Inches(2.2), Inches(11), Inches(2.5),
             "Корпоративное хранилище\nданных футбольной\nстатистики",
             size=48, color=TXT_PRIMARY, bold=True)

    add_text(slide, Inches(1.2), Inches(4.7), Inches(11), Inches(0.5),
             "на основе методологии Data Vault 2.0",
             size=22, color=TXT_MUTED)

    # Нижний блок автора
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(5.7), Inches(2), Emu(38000))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    add_text(slide, Inches(1.2), Inches(5.9), Inches(11), Inches(0.5),
             "Автор: Семёнов Антон  ·  3 курс  ·  2026",
             size=18, color=TXT_PRIMARY)

    add_text(slide, Inches(1.2), Inches(6.4), Inches(11), Inches(0.5),
             "Стек: Airflow · MinIO · PostgreSQL · Spark · ClickHouse · Superset",
             size=14, color=TXT_MUTED)


def slide_problem(prs: Presentation) -> None:
    slide = make_blank_slide(prs)
    slide_header(slide, 1, "Проблема и цель")

    add_bullets(slide, Inches(0.5), Inches(2.0), Inches(12), Inches(4.5), [
        "Спортивная аналитика — растущая индустрия; xG стал стандартной метрикой",
        "Открытые источники (Understat, StatsBomb) дают сырые данные, но без структуры",
        "Коммерческие платформы (FBref, Wyscout, StatsBomb IQ) — закрытый код, дорогие лицензии",
    ], size=20)

    # Цель — в виде акцентной плашки
    target = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.6))
    target.fill.solid()
    target.fill.fore_color.rgb = BG_PANEL
    target.line.color.rgb = ACCENT
    target.line.width = Pt(2)
    target.shadow.inherit = False

    add_text(slide, Inches(0.8), Inches(5.45), Inches(12), Inches(0.5),
             "ЦЕЛЬ", size=12, color=ACCENT, bold=True)
    add_text(slide, Inches(0.8), Inches(5.85), Inches(12), Inches(1.0),
             "Спроектировать DWH для футбольной статистики, продемонстрировать "
             "применение Data Vault 2.0 на современном open-source стеке.",
             size=16, color=TXT_PRIMARY)


def slide_sources(prs: Presentation) -> None:
    slide = make_blank_slide(prs)
    slide_header(slide, 2, "Источники данных")

    # Левая карточка — Understat
    card1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(0.5), Inches(2.2), Inches(6.0), Inches(4.7))
    card1.fill.solid()
    card1.fill.fore_color.rgb = BG_PANEL
    card1.line.color.rgb = ACCENT
    card1.line.width = Pt(1)
    card1.shadow.inherit = False

    add_text(slide, Inches(0.85), Inches(2.45), Inches(5.5), Inches(0.5),
             "UNDERSTAT.COM", size=14, color=ACCENT, bold=True)
    add_text(slide, Inches(0.85), Inches(2.95), Inches(5.5), Inches(0.8),
             "Основной источник", size=22, color=TXT_PRIMARY, bold=True)
    add_bullets(slide, Inches(0.85), Inches(4.0), Inches(5.5), Inches(2.7), [
        "HTML-страницы с встроенным JSON",
        "Топ-5 лиг Европы, сезоны 2014+",
        "Метрики: xG, xA, xPTS, PPDA",
        "Расширенная разметка матчей",
    ], size=14)

    # Правая карточка — StatsBomb
    card2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(6.85), Inches(2.2), Inches(6.0), Inches(4.7))
    card2.fill.solid()
    card2.fill.fore_color.rgb = BG_PANEL
    card2.line.color.rgb = ACCENT_2
    card2.line.width = Pt(1)
    card2.shadow.inherit = False

    add_text(slide, Inches(7.2), Inches(2.45), Inches(5.5), Inches(0.5),
             "STATSBOMB OPEN DATA", size=14, color=ACCENT_2, bold=True)
    add_text(slide, Inches(7.2), Inches(2.95), Inches(5.5), Inches(0.8),
             "Дополнительный источник", size=22, color=TXT_PRIMARY, bold=True)
    add_bullets(slide, Inches(7.2), Inches(4.0), Inches(5.5), Inches(2.7), [
        "JSON в публичном GitHub-репо",
        "La Liga (Barcelona-фокус, 2004-2020)",
        "Полный сезон La Liga 2015/16",
        "Детальная событийная разметка",
    ], size=14)


def slide_stack(prs: Presentation) -> None:
    slide = make_blank_slide(prs)
    slide_header(slide, 3, "Стек технологий")

    # 6 карточек технологий — 2 ряда по 3
    items = [
        ("Apache Airflow 2.9", "Оркестрация DAG-ов\nDatasets-цепочка", ACCENT),
        ("MinIO", "S3-совместимый\nraw lake", ACCENT_3),
        ("PostgreSQL 16", "OLTP с DV-схемой\n+ marts", ACCENT_2),
        ("dbt + datavault4dbt", "Сборка RV/BV/Marts\nиз макросов", ACCENT),
        ("Apache Spark 3.5", "Расчёт Elo\nперелив витрин", ACCENT_3),
        ("ClickHouse + Superset", "OLAP + BI\n2 дашборда", ACCENT_2),
    ]
    cw = Inches(4.0)
    ch = Inches(2.0)
    gap = Inches(0.15)
    start_left = Inches(0.5)
    start_top = Inches(2.2)

    for i, (title, desc, color) in enumerate(items):
        col = i % 3
        row = i // 3
        x = start_left + col * (cw + gap)
        y = start_top + row * (ch + gap)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, ch)
        card.fill.solid()
        card.fill.fore_color.rgb = BG_PANEL
        card.line.color.rgb = color
        card.line.width = Pt(1.5)
        card.shadow.inherit = False
        add_text(slide, x + Inches(0.3), y + Inches(0.25), cw - Inches(0.6), Inches(0.6),
                 title, size=16, color=color, bold=True)
        add_text(slide, x + Inches(0.3), y + Inches(0.95), cw - Inches(0.6), Inches(1.0),
                 desc, size=12, color=TXT_PRIMARY)


def slide_diagram(prs: Presentation, num: int, title: str, image_path: Path,
                  caption: str) -> None:
    slide = make_blank_slide(prs)
    slide_header(slide, num, title)

    if image_path.exists():
        # Подгоняем картинку пропорционально под доступную область,
        # не растягивая по обеим осям (чтобы не было искажений)
        from PIL import Image
        with Image.open(image_path) as img:
            iw, ih = img.size
        avail_w = Inches(12.0)
        avail_h = Inches(4.85)
        ratio_w = avail_w / iw
        ratio_h = avail_h / ih
        ratio = min(ratio_w, ratio_h)
        pic_w = int(iw * ratio)
        pic_h = int(ih * ratio)
        # центрируем
        left = (SLIDE_W - pic_w) // 2
        top = Inches(1.95) + (avail_h - pic_h) // 2
        slide.shapes.add_picture(str(image_path), left, top, width=pic_w, height=pic_h)
    else:
        add_text(slide, Inches(0.5), Inches(3.5), Inches(12), Inches(0.5),
                 f"[image not found: {image_path.name}]", size=14, color=ACCENT_3,
                 align=PP_ALIGN.CENTER)

    add_text(slide, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.4),
             caption, size=12, color=TXT_MUTED, align=PP_ALIGN.CENTER)


def slide_dv_intro(prs: Presentation) -> None:
    slide = make_blank_slide(prs)
    slide_header(slide, 7, "Data Vault 2.0: почему именно он?")

    reasons = [
        ("MULTI-SOURCE БЕЗ БОЛИ",
         "Один hub_match для SB и Understat через префикс BK"),
        ("ПОЛНЫЙ АУДИТ",
         "Каждое изменение — новая строка satellite (ldts + hashdiff)"),
        ("ИДЕМПОТЕНТНОСТЬ",
         "Повторная загрузка не плодит дубли — hashdiff фильтрует"),
        ("РАСШИРЯЕМОСТЬ",
         "Третий источник = новый stage. RV не меняется"),
    ]
    cw = Inches(6.1)
    ch = Inches(2.1)
    gap = Inches(0.2)
    for i, (title, desc) in enumerate(reasons):
        col = i % 2
        row = i // 2
        x = Inches(0.5) + col * (cw + gap)
        y = Inches(2.0) + row * (ch + gap)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, ch)
        card.fill.solid()
        card.fill.fore_color.rgb = BG_PANEL
        card.line.color.rgb = ACCENT
        card.line.width = Pt(1)
        card.shadow.inherit = False
        # большой номер
        add_text(slide, x + Inches(0.3), y + Inches(0.2), Inches(0.8), Inches(0.7),
                 f"0{i+1}", size=28, color=ACCENT, bold=True)
        add_text(slide, x + Inches(1.1), y + Inches(0.3), cw - Inches(1.4), Inches(0.6),
                 title, size=14, color=ACCENT, bold=True)
        add_text(slide, x + Inches(1.1), y + Inches(0.85), cw - Inches(1.4), Inches(1.2),
                 desc, size=14, color=TXT_PRIMARY)


def slide_marts_elo(prs: Presentation) -> None:
    slide = make_blank_slide(prs)
    slide_header(slide, 9, "Витрины и Spark Elo")

    # Левая колонка — список витрин
    add_text(slide, Inches(0.5), Inches(2.0), Inches(6), Inches(0.5),
             "8 АНАЛИТИЧЕСКИХ ВИТРИН", size=14, color=ACCENT, bold=True)

    marts = [
        ("mart_league_table", "386"),
        ("mart_top_scorers", "11 070"),
        ("mart_match_facts", "6 943"),
        ("mart_player_overperformers", "7 832"),
        ("mart_team_xg_trend", "386"),
        ("mart_sb_la_liga_history", "18"),
        ("mart_team_elo_current", "125"),
        ("mart_team_elo_history", "13 802"),
    ]
    for i, (name, count) in enumerate(marts):
        y = Inches(2.55) + i * Inches(0.42)
        add_text(slide, Inches(0.5), y, Inches(4.5), Inches(0.4),
                 f"  {name}", size=13, color=TXT_PRIMARY)
        add_text(slide, Inches(5.0), y, Inches(1.5), Inches(0.4),
                 count, size=13, color=ACCENT_2, bold=True, align=PP_ALIGN.RIGHT)

    # Правая колонка — Elo
    elo_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(7.0), Inches(2.0), Inches(5.85), Inches(4.9))
    elo_card.fill.solid()
    elo_card.fill.fore_color.rgb = BG_PANEL
    elo_card.line.color.rgb = ACCENT_3
    elo_card.line.width = Pt(1.5)
    elo_card.shadow.inherit = False

    add_text(slide, Inches(7.3), Inches(2.2), Inches(5.5), Inches(0.5),
             "SPARK: РАСЧЁТ ELO", size=14, color=ACCENT_3, bold=True)
    add_text(slide, Inches(7.3), Inches(2.7), Inches(5.5), Inches(0.6),
             "Per-league ClubElo", size=22, color=TXT_PRIMARY, bold=True)
    add_bullets(slide, Inches(7.3), Inches(3.6), Inches(5.5), Inches(3.0), [
        "Стартовый рейтинг 1500",
        "K-фактор = 20, home advantage = +100",
        "GD-modifier: ln(|gd|+1) при |gd| ≥ 2",
        "Per-league пул (изолированные лиги)",
        "6 943 матча → расчёт < 1 секунды",
    ], size=14)


def slide_pg_to_ch(prs: Presentation) -> None:
    slide = make_blank_slide(prs)
    slide_header(slide, 10, "Перелив PostgreSQL → ClickHouse")

    add_text(slide, Inches(0.5), Inches(2.0), Inches(12), Inches(0.6),
             "Почему через Parquet, а не напрямую?", size=20, color=TXT_PRIMARY, bold=True)

    add_bullets(slide, Inches(0.5), Inches(2.7), Inches(12), Inches(2.0), [
        "ClickHouse JDBC-driver не входит в стандартный Spark образ",
        "Spark s3a-writer требует aws-java-sdk-bundle (~273 МБ через Maven)",
        "Maven Central из РФ нестабильно отвечает — connection refused",
    ], size=15)

    # Стрелочный pipeline
    boxes = ["PostgreSQL", "/tmp/parquet", "MinIO marts/", "ClickHouse"]
    actions = ["Spark JDBC.read", "mc cp", "INSERT FROM s3()"]
    bw = Inches(2.6)
    bh = Inches(1.0)
    gap = Inches(0.55)
    start_left = Inches(0.5)
    y = Inches(5.4)

    for i, name in enumerate(boxes):
        x = start_left + i * (bw + gap)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, bw, bh)
        box.fill.solid()
        box.fill.fore_color.rgb = BG_PANEL
        box.line.color.rgb = ACCENT
        box.line.width = Pt(1.5)
        box.shadow.inherit = False
        add_text(slide, x, y + Inches(0.3), bw, Inches(0.5),
                 name, size=14, color=TXT_PRIMARY, bold=True, align=PP_ALIGN.CENTER)
        if i < len(boxes) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                            x + bw + Inches(0.1), y + Inches(0.4), gap - Inches(0.2), Inches(0.2))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT
            arrow.line.fill.background()
            add_text(slide, x + bw, y + Inches(0.65), gap, Inches(0.4),
                     actions[i], size=10, color=TXT_MUTED, align=PP_ALIGN.CENTER)


def slide_dashboard(prs: Presentation, num: int, title: str, subtitle: str,
                    image_path: Path) -> None:
    slide = make_blank_slide(prs)
    slide_header(slide, num, title)
    add_text(slide, Inches(0.5), Inches(1.85), Inches(12), Inches(0.5),
             subtitle, size=15, color=TXT_MUTED)

    if image_path.exists():
        slide.shapes.add_picture(str(image_path),
                                  Inches(1.5), Inches(2.5),
                                  width=Inches(10.3))
    else:
        add_text(slide, Inches(0.5), Inches(3.5), Inches(12), Inches(0.5),
                 f"[image not found: {image_path.name}]", size=14, color=ACCENT_3,
                 align=PP_ALIGN.CENTER)


def slide_outro(prs: Presentation) -> None:
    slide = make_blank_slide(prs)
    slide_header(slide, 13, "Что я узнал и куда дальше")

    # Слева — освоено
    add_text(slide, Inches(0.5), Inches(2.0), Inches(6), Inches(0.5),
             "ОСВОЕНО", size=14, color=ACCENT_2, bold=True)
    add_bullets(slide, Inches(0.5), Inches(2.55), Inches(6), Inches(4), [
        "Data Vault 2.0 на практике",
        "Airflow Datasets-цепочка",
        "Spark JDBC + DataFrame write",
        "Superset Native Filters + REST API",
        "ClickHouse INSERT FROM s3()",
    ], size=15)

    # Справа — дальше
    add_text(slide, Inches(7.0), Inches(2.0), Inches(6), Inches(0.5),
             "ДАЛЬШЕ", size=14, color=ACCENT_3, bold=True)
    add_bullets(slide, Inches(7.0), Inches(2.55), Inches(6), Inches(4), [
        "CI: ruff + sqlfluff + dbt test",
        "Инкрементальные модели (incremental)",
        "ML-надстройка: свой xG-классификатор",
        "Real-time event-stream от StatsBomb",
        "Replace docker exec на SparkSubmitOperator",
    ], size=15)


def slide_thanks(prs: Presentation) -> None:
    slide = make_blank_slide(prs)

    side = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.4), SLIDE_H)
    side.fill.solid()
    side.fill.fore_color.rgb = ACCENT
    side.line.fill.background()

    add_text(slide, Inches(1.2), Inches(2.8), Inches(11), Inches(2.0),
             "Спасибо\nза внимание!", size=64, color=TXT_PRIMARY, bold=True)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(1.2), Inches(5.2), Inches(2), Emu(38000))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    add_text(slide, Inches(1.2), Inches(5.4), Inches(11), Inches(0.6),
             "Q&A", size=28, color=ACCENT, bold=True)
    add_text(slide, Inches(1.2), Inches(6.0), Inches(11), Inches(0.5),
             "github.com/asemenov2004/football-dwh",
             size=14, color=TXT_MUTED)


def main() -> int:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    print("building presentation...")
    slide_title(prs)
    print("  01: title")
    slide_problem(prs)
    print("  02: problem")
    slide_sources(prs)
    print("  03: sources")
    slide_stack(prs)
    print("  04: stack")
    slide_diagram(prs, 5, "Архитектура: C4 Context",
                   PNG / "01_c4_context.png",
                   "Внешние границы Football DWH: один пользователь и два источника")
    print("  05: c4 context")
    slide_diagram(prs, 6, "Архитектура: C4 Containers",
                   PNG / "02_c4_containers.png",
                   "6 контейнеров и протоколы между ними")
    print("  06: c4 containers")
    slide_diagram(prs, 7, "Поток данных (DFD)",
                   PNG / "03_dfd.png",
                   "8 шагов от ingestion до Superset-чартов")
    print("  07: dfd")
    slide_dv_intro(prs)
    print("  08: dv intro")
    slide_diagram(prs, 9, "Data Vault 2.0: Hubs / Links / Satellites",
                   PNG / "04_er_data_vault_1.png",
                   "5 хабов · 5 линков · 6 сателлитов · MD5-hashing")
    print("  09: dv er")
    slide_marts_elo(prs)
    print("  10: marts + elo")
    slide_pg_to_ch(prs)
    print("  11: pg to ch")
    slide_dashboard(prs, 12, "Дашборд 1: Football DWH",
                     "9 чартов с фильтрами по лиге и сезону",
                     SCR / "01_football_dwh.png")
    print("  12: dashboard 1")
    slide_dashboard(prs, 13, "Дашборд 2: European Teams",
                     "4 кросс-лиговых чарта без фильтра лиги",
                     SCR / "02_european_teams.png")
    print("  13: dashboard 2")
    slide_outro(prs)
    print("  14: outro")
    slide_thanks(prs)
    print("  15: thanks")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"\nDONE. {OUT} ({OUT.stat().st_size // 1024} KiB, {len(prs.slides)} слайдов)")
    return 0


if __name__ == "__main__":
    sys.exit(main())
